import os
import pandas as pd
from dataclasses import dataclass, field
from typing import Optional, List, Dict, Any


@dataclass
class FileData:
    filename: str
    file_type: str  # 'tabular' | 'document'
    dataframe: Optional[pd.DataFrame] = None
    sheet_names: Optional[List[str]] = None
    dataframes: Optional[Dict[str, pd.DataFrame]] = None
    text_content: Optional[str] = None
    tables: Optional[List[pd.DataFrame]] = None
    metadata: Optional[Dict[str, Any]] = None
    slides: Optional[List[Dict]] = None


_PROCESSORS: Dict[str, Any] = {}


def process_file(file_path: str) -> FileData:
    filename = os.path.basename(file_path)
    ext = os.path.splitext(filename)[1].lower()
    handlers = {
        '.csv':  _csv,
        '.xlsx': _excel,
        '.xls':  _excel,
        '.pdf':  _pdf,
        '.docx': _docx,
        '.pptx': _pptx,
    }
    if ext not in handlers:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {', '.join(handlers)}")
    return handlers[ext](file_path, filename)


# ── Tabular ──────────────────────────────────────────────────────────────────

def _csv(path: str, filename: str) -> FileData:
    df = pd.read_csv(path)
    return FileData(filename=filename, file_type='tabular', dataframe=df)


def _excel(path: str, filename: str) -> FileData:
    xl = pd.ExcelFile(path)
    sheets = xl.sheet_names
    dfs = {s: xl.parse(s) for s in sheets}
    return FileData(
        filename=filename,
        file_type='tabular',
        dataframe=dfs[sheets[0]],
        sheet_names=sheets,
        dataframes=dfs,
    )


# ── Documents ─────────────────────────────────────────────────────────────────

def _pdf(path: str, filename: str) -> FileData:
    import pdfplumber
    texts, tables = [], []
    with pdfplumber.open(path) as pdf:
        meta: Dict[str, Any] = {'page_count': len(pdf.pages)}
        if pdf.metadata:
            meta.update({
                k: v for k, v in pdf.metadata.items()
                if v and isinstance(v, (str, int, float))
            })
        for i, page in enumerate(pdf.pages, 1):
            t = page.extract_text()
            if t:
                texts.append(f"[Page {i}]\n{t}")
            for tbl in page.extract_tables() or []:
                if len(tbl) > 1:
                    try:
                        df = pd.DataFrame(tbl[1:], columns=tbl[0])
                        df.attrs['source_page'] = i
                        tables.append(df)
                    except Exception:
                        pass
    return FileData(
        filename=filename,
        file_type='document',
        text_content='\n\n'.join(texts),
        tables=tables or None,
        metadata=meta,
    )


def _docx(path: str, filename: str) -> FileData:
    from docx import Document
    doc = Document(path)
    props = doc.core_properties
    meta: Dict[str, Any] = {
        'author': props.author,
        'created': str(props.created) if props.created else None,
        'word_count': sum(len(p.text.split()) for p in doc.paragraphs),
    }
    texts = [p.text for p in doc.paragraphs if p.text.strip()]
    tables = []
    for tbl in doc.tables:
        rows = [[c.text for c in row.cells] for row in tbl.rows]
        if len(rows) > 1:
            try:
                tables.append(pd.DataFrame(rows[1:], columns=rows[0]))
            except Exception:
                pass
    return FileData(
        filename=filename,
        file_type='document',
        text_content='\n'.join(texts),
        tables=tables or None,
        metadata=meta,
    )


def _pptx(path: str, filename: str) -> FileData:
    from pptx import Presentation
    prs = Presentation(path)
    slides_data: List[Dict] = []
    texts: List[str] = []

    for i, slide in enumerate(prs.slides, 1):
        info: Dict[str, Any] = {
            'slide_number': i, 'title': '', 'content': [], 'notes': ''
        }
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if not t:
                continue
            pf = getattr(shape, 'placeholder_format', None)
            if pf is not None and pf.idx == 0:
                info['title'] = t
            else:
                info['content'].append(t)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                info['notes'] = notes

        slides_data.append(info)
        body = '\n'.join(info['content'])
        slide_text = f"[Slide {i}: {info['title']}]\n{body}"
        if info['notes']:
            slide_text += f"\n[Notes: {info['notes']}]"
        texts.append(slide_text)

    return FileData(
        filename=filename,
        file_type='document',
        text_content='\n\n'.join(texts),
        slides=slides_data,
        metadata={'slide_count': len(prs.slides)},
    )
